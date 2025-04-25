import streamlit as st
from pathlib import Path
import fitz  # PyMuPDF

# Function to extract text from uploaded file
def extract_text_from_file(uploaded_file):
    file_extension = Path(uploaded_file.name).suffix.lower()
    
    if file_extension == ".pdf":
        with fitz.open(stream=uploaded_file.read(), filetype="pdf") as doc:
            text = ""
            for page in doc:
                text += page.get_text()
            return text if text else "No text found in PDF."
    
    elif file_extension == ".docx":
        from docx import Document
        doc = Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    
    elif file_extension == ".pptx":
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text if text else "No text found in PPT."
    
    else:
        return "Unsupported file format."

# Dummy data for query output
dummy_results = [
    {"title": "Document A", "summary": "This is a summary of Document A.", "related": ["Document C", "Document D"]},
    {"title": "Document B", "summary": "This is a summary of Document B.", "related": ["Document A"]}
]

# Streamlit App
st.title("AI-Based Document Search Assistant")

# Upload files
st.header("1. Upload Your Documents")
uploaded_files = st.file_uploader("Upload PDF, Word, or PPT files", type=["pdf", "docx", "pptx"], accept_multiple_files=True)

if uploaded_files:
    st.success(f"{len(uploaded_files)} file(s) uploaded successfully.")
    for file in uploaded_files:
        st.write(f"• {file.name}")
        with st.expander(f"Preview: {file.name}"):
            extracted_text = extract_text_from_file(file)
            st.text_area("Extracted Text", extracted_text[:2000], height=200, key=file.name)

# Query input
st.header("2. Enter Your Search Query")
query = st.text_input("Type your question or query here")

# Display results
if query:
    st.header("3. Search Results")
    st.info(f"Showing results for: **{query}**")

    for i, result in enumerate(dummy_results, 1):
        st.subheader(f"{i}. {result['title']}")
        st.write(f"**Summary:** {result['summary']}")
        st.write("**Related Documents:**")
        for rel in result["related"]:
            st.write(f"- {rel}")

# Footer
st.markdown("---")
st.caption("Prototype Interface | Team Project")
