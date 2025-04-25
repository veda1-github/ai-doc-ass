from flask import Flask, request, jsonify
from werkzeug.utils import secure_filename
import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer, util
import numpy as np
from typing import Dict, List
import uuid
from transformers import pipeline
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)

# Configuration
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Initialize NLP models with error handling
try:
    # Try to load summarization model (smaller version if the large one fails)
    try:
        summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
        logger.info("Loaded facebook/bart-large-cnn model")
    except Exception as e:
        logger.warning(f"Could not load bart-large-cnn, trying smaller model: {str(e)}")
        summarizer = pipeline("summarization", model="facebook/bart-base")
        logger.info("Loaded facebook/bart-base model as fallback")
    
    # Load embedding model
    embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
    logger.info("Loaded sentence transformer model")
except Exception as e:
    logger.error(f"Failed to load models: {str(e)}")
    raise

# In-memory storage for documents
documents_db: Dict[str, Dict] = {}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text(filepath: str, filename: str) -> str:
    """Extract text from different file formats"""
    ext = filename.rsplit('.', 1)[1].lower()
    
    try:
        if ext == 'pdf':
            with fitz.open(filepath) as doc:
                return "\n".join([page.get_text() for page in doc])
        elif ext == 'docx':
            doc = Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        elif ext == 'pptx':
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
    except Exception as e:
        logger.error(f"Error extracting text from {filename}: {str(e)}")
        return ""

def summarize_text(text: str) -> str:
    """Summarize text using NLP model with chunking"""
    if not text.strip():
        return ""
    
    try:
        # Split into chunks if too long
        max_chunk_size = 1024  # tokens
        words = text.split()
        chunks = [' '.join(words[i:i+max_chunk_size]) for i in range(0, len(words), max_chunk_size)]
        
        summaries = []
        for chunk in chunks:
            summary = summarizer(chunk, max_length=130, min_length=30, do_sample=False)[0]['summary_text']
            summaries.append(summary)
        
        return " ".join(summaries)
    except Exception as e:
        logger.error(f"Error summarizing text: {str(e)}")
        # Fallback to first 500 characters if summarization fails
        return text[:500]

def process_document(filepath: str, filename: str) -> Dict:
    """Process an uploaded document"""
    full_text = extract_text(filepath, filename)
    if not full_text:
        return None
    
    summary = summarize_text(full_text)
    
    try:
        embeddings = embedding_model.encode(full_text, convert_to_tensor=True)
        return {
            'filename': filename,
            'full_text': full_text,
            'summary': summary,
            'embeddings': embeddings.cpu().numpy()
        }
    except Exception as e:
        logger.error(f"Error creating embeddings for {filename}: {str(e)}")
        return None

@app.route('/upload', methods=['POST'])
def upload_files():
    """Handle file upload and processing"""
    if 'files' not in request.files:
        return jsonify({'error': 'No files provided'}), 400
    
    files = request.files.getlist('files')
    results = []
    
    for file in files:
        if file.filename == '':
            continue
            
        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            file_id = str(uuid.uuid4())
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
            file.save(filepath)
            
            doc = process_document(filepath, filename)
            if doc:
                documents_db[file_id] = doc
                results.append({
                    'id': file_id,
                    'filename': filename,
                    'status': 'processed',
                    'summary': doc['summary']
                })
            else:
                results.append({
                    'filename': filename,
                    'status': 'error',
                    'message': 'Could not process file'
                })
            
            os.remove(filepath)
    
    return jsonify({'results': results})

@app.route('/search', methods=['POST'])
def search_documents():
    """Search across processed documents"""
    data = request.get_json()
    query = data.get('query', '').strip()
    
    if not query:
        return jsonify({'error': 'Empty query'}), 400
    
    if not documents_db:
        return jsonify({'error': 'No documents processed'}), 400
    
    try:
        query_embedding = embedding_model.encode(query, convert_to_tensor=True)
        results = []
        
        for doc_id, doc in documents_db.items():
            doc_embedding = doc['embeddings']
            similarity = util.pytorch_cos_sim(query_embedding, doc_embedding).item()
            
            sentences = [s.strip() for s in doc['full_text'].split('.') if s.strip()]
            if sentences:
                sentence_embeddings = embedding_model.encode(sentences, convert_to_tensor=True)
                similarities = util.pytorch_cos_sim(query_embedding, sentence_embeddings)[0]
                top_indices = similarities.argsort(descending=True)[:3]
                relevant_sentences = [sentences[i] for i in top_indices if i < len(sentences)]
            else:
                relevant_sentences = ["No text content found"]
            
            results.append({
                'file': doc['filename'],
                'content': ' '.join(relevant_sentences),
                'score': similarity
            })
        
        results.sort(key=lambda x: x['score'], reverse=True)
        return jsonify(results[:10])  # Return top 10 results
    
    except Exception as e:
        logger.error(f"Search error: {str(e)}")
        return jsonify({'error': 'Search failed'}), 500

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({'status': 'healthy', 'model': 'loaded'})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)